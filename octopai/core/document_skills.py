"""
Octopai Document Skills - Advanced Document Processing Capabilities

This module provides comprehensive document processing skills for Octopai,
enabling AI Agents to work with various document formats including:
- PDF documents (extraction, creation, manipulation)
- Word documents (DOCX) - reading, editing, creation
- Excel spreadsheets (XLSX) - data analysis, charting
- PowerPoint presentations (PPTX) - slide manipulation
- And more document formats

These skills form Octopai's document intelligence layer,
empowering AI Agents with sophisticated document processing capabilities.
"""

import os
from typing import Dict, List, Any, Optional, Union, BinaryIO
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from datetime import datetime


class DocumentFormat(Enum):
    """Supported document formats"""
    PDF = "pdf"
    DOCX = "docx"
    XLSX = "xlsx"
    PPTX = "pptx"
    TXT = "txt"
    MD = "markdown"
    CSV = "csv"
    JSON = "json"
    HTML = "html"


@dataclass
class DocumentMetadata:
    """Metadata for a document"""
    title: Optional[str] = None
    author: Optional[str] = None
    subject: Optional[str] = None
    keywords: List[str] = field(default_factory=list)
    created_at: Optional[str] = None
    modified_at: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    file_size: Optional[int] = None
    custom_fields: Dict[str, Any] = field(default_factory=dict)
    
    def to_dict(self) -> Dict[str, Any]:
        return {
            "title": self.title,
            "author": self.author,
            "subject": self.subject,
            "keywords": self.keywords,
            "created_at": self.created_at,
            "modified_at": self.modified_at,
            "page_count": self.page_count,
            "word_count": self.word_count,
            "file_size": self.file_size,
            "custom_fields": self.custom_fields
        }


@dataclass
class ExtractedText:
    """Result of text extraction from a document"""
    text: str
    pages: List[str] = field(default_factory=list)
    metadata: Optional[DocumentMetadata] = None
    sections: List[Dict[str, Any]] = field(default_factory=list)


@dataclass
class FormField:
    """Form field in a document"""
    name: str
    value: Optional[str] = None
    field_type: str = "text"
    options: List[str] = field(default_factory=list)
    required: bool = False


@dataclass
class TableData:
    """Table data extracted from a document"""
    table_id: str
    headers: List[str]
    rows: List[List[Any]]
    sheet_name: Optional[str] = None


class PDFSkill:
    """
    Octopai PDF Skill - Comprehensive PDF Processing
    
    Provides capabilities for:
    - Text extraction from PDFs
    - Form field extraction and manipulation
    - Metadata reading/writing
    - Page manipulation
    - PDF creation and merging
    """
    
    @staticmethod
    def extract_text(
        pdf_path: Union[str, Path],
        page_numbers: Optional[List[int]] = None
    ) -> ExtractedText:
        """
        Extract text from a PDF document
        
        Args:
            pdf_path: Path to the PDF file
            page_numbers: Optional list of page numbers to extract (1-based)
            
        Returns:
            ExtractedText object with text content and metadata
        """
        try:
            import PyPDF2
            
            pdf_path = Path(pdf_path)
            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                
                metadata = DocumentMetadata(
                    title=reader.metadata.get('/Title'),
                    author=reader.metadata.get('/Author'),
                    subject=reader.metadata.get('/Subject'),
                    page_count=len(reader.pages),
                    file_size=pdf_path.stat().st_size
                )
                
                pages = []
                all_text = []
                
                for i, page in enumerate(reader.pages):
                    if page_numbers and (i + 1) not in page_numbers:
                        continue
                    text = page.extract_text() or ""
                    pages.append(text)
                    all_text.append(text)
                
                return ExtractedText(
                    text='\n\n'.join(all_text),
                    pages=pages,
                    metadata=metadata
                )
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF processing. Install it with 'pip install PyPDF2'")
    
    @staticmethod
    def extract_form_fields(pdf_path: Union[str, Path]) -> List[FormField]:
        """
        Extract form fields from a PDF
        
        Args:
            pdf_path: Path to the PDF file
            
        Returns:
            List of FormField objects
        """
        try:
            import PyPDF2
            
            fields = []
            with open(pdf_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if '/AcroForm' in reader.trailer['/Root']:
                    form = reader.trailer['/Root']['/AcroForm']
                    if '/Fields' in form:
                        for field in form['/Fields']:
                            field_obj = field.get_object()
                            fields.append(FormField(
                                name=field_obj.get('/T', ''),
                                value=field_obj.get('/V'),
                                field_type=field_obj.get('/FT', 'text')
                            ))
            
            return fields
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF form processing")
    
    @staticmethod
    def merge_pdfs(
        input_paths: List[Union[str, Path]],
        output_path: Union[str, Path]
    ) -> bool:
        """
        Merge multiple PDF files into one
        
        Args:
            input_paths: List of PDF file paths to merge
            output_path: Output PDF file path
            
        Returns:
            True if successful
        """
        try:
            import PyPDF2
            
            merger = PyPDF2.PdfMerger()
            
            for path in input_paths:
                merger.append(str(path))
            
            with open(output_path, 'wb') as f:
                merger.write(f)
            
            merger.close()
            return True
        except ImportError:
            raise ImportError("PyPDF2 is required for PDF merging")


class DOCXSkill:
    """
    Octopai DOCX Skill - Word Document Processing
    
    Provides capabilities for:
    - Reading and writing Word documents
    - Text extraction and manipulation
    - Style and formatting management
    - Table creation and manipulation
    - Document creation from templates
    """
    
    @staticmethod
    def extract_text(docx_path: Union[str, Path]) -> ExtractedText:
        """
        Extract text from a Word document
        
        Args:
            docx_path: Path to the DOCX file
            
        Returns:
            ExtractedText object with text content and metadata
        """
        try:
            from docx import Document
            
            doc = Document(docx_path)
            docx_path = Path(docx_path)
            
            paragraphs = [p.text for p in doc.paragraphs]
            text = '\n\n'.join(paragraphs)
            
            metadata = DocumentMetadata(
                title=doc.core_properties.title,
                author=doc.core_properties.author,
                subject=doc.core_properties.subject,
                keywords=doc.core_properties.keywords.split(',') if doc.core_properties.keywords else [],
                created_at=doc.core_properties.created.isoformat() if doc.core_properties.created else None,
                modified_at=doc.core_properties.modified.isoformat() if doc.core_properties.modified else None,
                file_size=docx_path.stat().st_size
            )
            
            return ExtractedText(
                text=text,
                metadata=metadata,
                sections=[{"text": p.text, "style": p.style.name} for p in doc.paragraphs]
            )
        except ImportError:
            raise ImportError("python-docx is required for DOCX processing. Install it with 'pip install python-docx'")
    
    @staticmethod
    def extract_tables(docx_path: Union[str, Path]) -> List[TableData]:
        """
        Extract tables from a Word document
        
        Args:
            docx_path: Path to the DOCX file
            
        Returns:
            List of TableData objects
        """
        try:
            from docx import Document
            
            doc = Document(docx_path)
            tables = []
            
            for table_idx, table in enumerate(doc.tables):
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                rows = []
                
                for row in table.rows[1:]:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows.append(row_data)
                
                tables.append(TableData(
                    table_id=f"table_{table_idx}",
                    headers=headers,
                    rows=rows
                ))
            
            return tables
        except ImportError:
            raise ImportError("python-docx is required for DOCX table extraction")
    
    @staticmethod
    def create_document(
        output_path: Union[str, Path],
        content: str,
        title: Optional[str] = None,
        author: Optional[str] = None
    ) -> bool:
        """
        Create a new Word document
        
        Args:
            output_path: Output DOCX file path
            content: Text content for the document
            title: Optional document title
            author: Optional author name
            
        Returns:
            True if successful
        """
        try:
            from docx import Document
            
            doc = Document()
            
            if title:
                doc.add_heading(title, 0)
            
            for paragraph in content.split('\n\n'):
                if paragraph.strip():
                    doc.add_paragraph(paragraph.strip())
            
            if author:
                doc.core_properties.author = author
            
            doc.save(output_path)
            return True
        except ImportError:
            raise ImportError("python-docx is required for DOCX creation")


class XLSXSkill:
    """
    Octopai XLSX Skill - Spreadsheet Processing
    
    Provides capabilities for:
    - Reading and writing Excel spreadsheets
    - Data extraction and manipulation
    - Chart creation and manipulation
    - Formula evaluation
    - Pivot table creation
    """
    
    @staticmethod
    def read_sheet(
        xlsx_path: Union[str, Path],
        sheet_name: Optional[str] = None
    ) -> TableData:
        """
        Read a sheet from an Excel file
        
        Args:
            xlsx_path: Path to the XLSX file
            sheet_name: Optional sheet name (reads first sheet if not specified)
            
        Returns:
            TableData object with sheet content
        """
        try:
            import pandas as pd
            
            df = pd.read_excel(xlsx_path, sheet_name=sheet_name)
            
            headers = df.columns.tolist()
            rows = df.values.tolist()
            
            return TableData(
                table_id=sheet_name or "sheet_1",
                headers=headers,
                rows=rows,
                sheet_name=sheet_name
            )
        except ImportError:
            raise ImportError("pandas and openpyxl are required for XLSX processing. Install them with 'pip install pandas openpyxl'")
    
    @staticmethod
    def list_sheets(xlsx_path: Union[str, Path]) -> List[str]:
        """
        List all sheet names in an Excel file
        
        Args:
            xlsx_path: Path to the XLSX file
            
        Returns:
            List of sheet names
        """
        try:
            import pandas as pd
            
            excel_file = pd.ExcelFile(xlsx_path)
            return excel_file.sheet_names
        except ImportError:
            raise ImportError("pandas is required for XLSX sheet listing")
    
    @staticmethod
    def write_sheet(
        data: TableData,
        output_path: Union[str, Path],
        sheet_name: str = "Sheet1"
    ) -> bool:
        """
        Write data to an Excel sheet
        
        Args:
            data: TableData to write
            output_path: Output XLSX file path
            sheet_name: Sheet name
            
        Returns:
            True if successful
        """
        try:
            import pandas as pd
            
            df = pd.DataFrame(data.rows, columns=data.headers)
            
            mode = 'a' if Path(output_path).exists() else 'w'
            if_sheet_exists = 'replace' if mode == 'a' else None
            
            with pd.ExcelWriter(
                output_path,
                engine='openpyxl',
                mode=mode,
                if_sheet_exists=if_sheet_exists
            ) as writer:
                df.to_excel(writer, sheet_name=sheet_name, index=False)
            
            return True
        except ImportError:
            raise ImportError("pandas and openpyxl are required for XLSX writing")


class PPTXSkill:
    """
    Octopai PPTX Skill - Presentation Processing
    
    Provides capabilities for:
    - Reading and writing PowerPoint presentations
    - Slide extraction and manipulation
    - Text extraction from slides
    - Image extraction
    - Template-based presentation creation
    """
    
    @staticmethod
    def extract_text(pptx_path: Union[str, Path]) -> ExtractedText:
        """
        Extract text from a PowerPoint presentation
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            ExtractedText object with text content and metadata
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            pptx_path = Path(pptx_path)
            
            slide_texts = []
            all_text = []
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content.append(shape.text)
                
                slide_text = '\n'.join(slide_content)
                slide_texts.append(slide_text)
                all_text.append(f"--- Slide {slide_idx + 1} ---\n{slide_text}")
            
            metadata = DocumentMetadata(
                title=prs.core_properties.title,
                author=prs.core_properties.author,
                subject=prs.core_properties.subject,
                page_count=len(prs.slides),
                file_size=pptx_path.stat().st_size
            )
            
            return ExtractedText(
                text='\n\n'.join(all_text),
                pages=slide_texts,
                metadata=metadata
            )
        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing. Install it with 'pip install python-pptx'")
    
    @staticmethod
    def get_slide_count(pptx_path: Union[str, Path]) -> int:
        """
        Get the number of slides in a presentation
        
        Args:
            pptx_path: Path to the PPTX file
            
        Returns:
            Number of slides
        """
        try:
            from pptx import Presentation
            
            prs = Presentation(pptx_path)
            return len(prs.slides)
        except ImportError:
            raise ImportError("python-pptx is required for PPTX processing")


class DocumentSkillFactory:
    """
    Octopai Document Skill Factory - Unified Document Processing Interface
    
    Provides a unified interface for working with various document formats,
    automatically selecting the appropriate skill based on file type.
    """
    
    @staticmethod
    def get_skill_for_format(format: DocumentFormat):
        """
        Get the appropriate skill for a document format
        
        Args:
            format: Document format enum
            
        Returns:
            Skill class for the format
        """
        skill_map = {
            DocumentFormat.PDF: PDFSkill,
            DocumentFormat.DOCX: DOCXSkill,
            DocumentFormat.XLSX: XLSXSkill,
            DocumentFormat.PPTX: PPTXSkill,
        }
        return skill_map.get(format)
    
    @staticmethod
    def detect_format(file_path: Union[str, Path]) -> Optional[DocumentFormat]:
        """
        Detect document format from file extension
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Detected DocumentFormat or None
        """
        ext = Path(file_path).suffix.lower()
        format_map = {
            '.pdf': DocumentFormat.PDF,
            '.docx': DocumentFormat.DOCX,
            '.xlsx': DocumentFormat.XLSX,
            '.pptx': DocumentFormat.PPTX,
            '.txt': DocumentFormat.TXT,
            '.md': DocumentFormat.MD,
            '.csv': DocumentFormat.CSV,
            '.json': DocumentFormat.JSON,
            '.html': DocumentFormat.HTML,
        }
        return format_map.get(ext)
    
    @staticmethod
    def extract_text(file_path: Union[str, Path]) -> Optional[ExtractedText]:
        """
        Extract text from any supported document format
        
        Args:
            file_path: Path to the document file
            
        Returns:
            ExtractedText or None if format not supported
        """
        format = DocumentSkillFactory.detect_format(file_path)
        if not format:
            return None
        
        skill = DocumentSkillFactory.get_skill_for_format(format)
        if not skill:
            return None
        
        if hasattr(skill, 'extract_text'):
            return skill.extract_text(file_path)
        
        return None
    
    @staticmethod
    def extract_tables(file_path: Union[str, Path]) -> List[TableData]:
        """
        Extract tables from any supported document format
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of TableData objects
        """
        format = DocumentSkillFactory.detect_format(file_path)
        if not format:
            return []
        
        skill = DocumentSkillFactory.get_skill_for_format(format)
        if not skill:
            return []
        
        if hasattr(skill, 'extract_tables'):
            return skill.extract_tables(file_path)
        
        return []


__all__ = [
    'DocumentFormat',
    'DocumentMetadata',
    'ExtractedText',
    'FormField',
    'TableData',
    'PDFSkill',
    'DOCXSkill',
    'XLSXSkill',
    'PPTXSkill',
    'DocumentSkillFactory'
]
